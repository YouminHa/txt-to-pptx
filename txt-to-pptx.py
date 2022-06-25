"""
    txt-to-pptx.py

 A tool converting text file to ppt file, to make a propresenter slides from plain text files.

 Requires python-pptx  (pip install python-pptx)
"""

__author__ = 'Youmin Ha <youmin78.ha@gmail.com>'
__version__ = '0.1.0'
__license__ = 'MIT'

import copy
import getopt
import os.path
import sys

import pptx

def print_usage():
  print("Usage: python3 txt-to-ppt.py -t <template .pptx file> -i <input txt file> -o <output .pptx file>")


def read_lyrics(input_txt_path):
  # read txt into lyrics
  lyrics = []
  trans_line = []

  f = open(input_txt_path, 'r')
  eof = False
  while not eof:
    l = f.readline()
    if not l:
      eof = True
    l = l.strip()
    if len(l) == 0:  # empty line
      if len(trans_line) > 0:  # at least 1 line exists
        lyrics.append(trans_line)
        trans_line = []
    else:
      trans_line.append(l)

  f.close()

  return lyrics


def convert(template_ppt_path, input_txt_path, output_ppt_path):
  p = pptx.Presentation(template_ppt_path)

  src_slide = p.slides[0]

  lyrics = read_lyrics(input_txt_path)

  for slide_idx in range(len(lyrics)):
    trans_line = lyrics[slide_idx]

    # copy slide if necessary
    while len(p.slides) <= slide_idx:
      slide_layout = src_slide.slide_layout  # copy src_slide's layout
      new_slide = p.slides.add_slide(slide_layout)
      for shape in src_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    slide = p.slides[slide_idx]

    # copy text into slide shapes
    for line_idx in range(min(len(trans_line), len(slide.shapes))):
      # DEBUG
      print("slide {}, text {} : {}".format(slide_idx, line_idx, trans_line[line_idx]))

      # put the text into shape.text_frame.paragraph.text
      shape = slide.shapes[line_idx]
      if shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = trans_line[line_idx]

  p.save(output_ppt_path)

if __name__ == '__main__':
  template_pptx_path = ''
  input_txt_path = ''
  output_pptx_path = ''

  # commandline options
  options, args = getopt.getopt(sys.argv[1:], 'i:o:t:')
  for op, val in options:
    if op == '-i':
      input_txt_path = val
    elif op == '-o':
      output_pptx_path = val
    elif op == '-t':
      template_pptx_path = val

  if os.path.isfile(template_pptx_path) and os.path.isfile(input_txt_path):
      convert(template_pptx_path, input_txt_path, output_pptx_path)
  else:
    print_usage()
    sys.exit(1)

